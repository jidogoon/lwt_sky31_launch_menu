from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.table import Table, _Cell
from internal.menu_data import MenuData
from internal.sky31_abc import SKY31ABC
from internal.util import Util


class SKY31PPTX(SKY31ABC):
    def __init__(self, pptx_file: str):
        super().__init__()
        self._file_to_open = pptx_file

    def _load_menus(self):
        tables = self._get_menu_tables()
        for table in tables:
            self._parse_table(table)
        self._clear_prev_month()

    def _get_menu_tables(self):
        tables = []
        prs = Presentation(self._file_to_open)
        for slide in prs.slides:
            for shape in slide.shapes:
                if MSO_SHAPE_TYPE.TABLE is not shape.shape_type:
                    continue
                tables.append(shape.table)
        return tables

    def _parse_table(self, table: Table):
        for r_index, row in enumerate(table.rows):
            if not row.cells[0].text.endswith('주차'):
                continue
            for c_index, cell in enumerate(row.cells):
                if self._is_ignore_cell(c_index, cell):
                    continue
                self._find_menu(c_index, cell, r_index, table)

    def _is_ignore_cell(self, c_index: int, cell: _Cell) -> bool:
        return c_index < 1 or self._is_holiday_cell(cell)

    def _find_menu(self, c_index: int, cell: _Cell, r_index: int, table: Table):
        for i in range(1, 4):
            current_cell = table.rows[r_index + i].cells[c_index]
            if None is Util.only_num(current_cell.text):
                continue
            menu_text = current_cell.text.replace('\x0b', '\n').strip()
            if '' is menu_text:
                continue
            nth_day = Util.only_num(cell.text)
            menu_text = menu_text.replace('\n', '').replace(')', ')\n').strip()
            menu_text_spt = self._reduce_malformed_menu(menu_text).split('\n')
            self._append_menu(menu_text_spt, nth_day)
            if len(menu_text_spt) > 3:
                self._append_menu(menu_text_spt, nth_day, moremenu=True)

    def _reduce_malformed_menu(self, menu_text: str) -> str:
        bracket_idxs = self._find_idxes_in_text(menu_text, '(')
        bracket_newline_idxs = self._find_idxes_in_text(menu_text, '\n(')
        if len(bracket_idxs) != len(bracket_newline_idxs):
            replaced_cnt = 0
            for idx in bracket_idxs:
                if menu_text[idx + replaced_cnt - 1] != '\n':
                    menu_text = menu_text[:idx + replaced_cnt] + '\n' + menu_text[idx + replaced_cnt:]
                    replaced_cnt += 1
        return menu_text

    def _find_idxes_in_text(self, text: str, word: str):
        all_positions = []
        next_pos = -1
        while True:
            next_pos = text.find(word, next_pos + 1)
            if next_pos < 0:
                break
            all_positions.append(next_pos)
        return all_positions

    def _append_menu(self, menu_text_spt: [], nth_day: int, moremenu: bool = False):
        offset = 2 if moremenu else 0
        price = Util.only_num(menu_text_spt[1 + offset])
        if None is price:
            return
        self._menus.append(MenuData(nth_day=nth_day,
                                    menu=menu_text_spt[0 + offset],
                                    price=price))

    def _is_holiday_cell(self, cell: _Cell):
        for prg in cell.text_frame.paragraphs:
            for run in prg.runs:
                if MSO_COLOR_TYPE.RGB is run.font.color.type:
                    return True
        return False
