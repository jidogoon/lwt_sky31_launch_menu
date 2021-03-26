import re
from dataclasses import dataclass
from datetime import date

from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.table import Table, _Cell


@dataclass()
class MenuData:
    nth_day: int
    menu: str
    price: int


class SKY31:
    def __init__(self):
        yyyymm = date.today().strftime("%Y%m")
        self._file_to_open = f'{yyyymm}.pptx'
        self._menus = []

    def get_today_menus(self) -> []:
        self._load_menus()
        nth_day = date.today().day
        today_menus = list(filter(lambda m: m.nth_day == nth_day, self._menus))
        return today_menus

    def get_monthly_menus(self):
        return self._menus

    def _load_menus(self):
        tables = self._get_menu_tables()
        for table in tables:
            self._parse_table(table)

    def _get_menu_tables(self):
        tables = []
        prs = Presentation(self._file_to_open)
        for slide in prs.slides:
            for shape in slide.shapes:
                if MSO_SHAPE_TYPE.TABLE is not shape.shape_type:
                    continue
                tables.append(shape.table)
        return tables

    def _parse_table(self, table: Table):
        for r_index, row in enumerate(table.rows):
            if not row.cells[0].text.endswith('주차'):
                continue
            for c_index, cell in enumerate(row.cells):
                if self._is_ignore_cell(c_index, cell):
                    continue
                self._find_menu(c_index, cell, r_index, table)

    def _find_menu(self, c_index: int, cell: _Cell, r_index: int, table: Table):
        for i in range(1, 4):
            menu_text = table.rows[r_index + i].cells[c_index].text.replace('\x0b', '\n').strip()
            if '' is menu_text:
                continue
            nth_day = self._only_num(cell.text)
            menu_text_spt = menu_text.split('\n')
            self._append_menu(menu_text_spt, nth_day)
            if len(menu_text_spt) > 3:
                self._append_menu(menu_text_spt, nth_day, moremenu=True)

    def _is_ignore_cell(self, c_index: int, cell: _Cell) -> bool:
        return c_index < 1 or self._is_holiday_cell(cell)

    def _append_menu(self, menu_text_spt: [], nth_day: int, moremenu: bool = False):
        offset = 3 if moremenu else 0
        self._menus.append(MenuData(nth_day=nth_day,
                                    menu=menu_text_spt[0 + offset],
                                    price=self._only_num(menu_text_spt[1 + offset])))

    def _is_holiday_cell(self, cell: _Cell):
        for prg in cell.text_frame.paragraphs:
            for run in prg.runs:
                if MSO_COLOR_TYPE.RGB is run.font.color.type:
                    return True
        return False

    def _only_num(self, text) -> int:
        return int(''.join(re.findall('\d+', text)))


if __name__ == '__main__':
    today_menu = SKY31().get_today_menus()
    print(date.today().strftime("%Y년 %m월 %d일 SKY 31 점심메뉴"))
    for menu in today_menu:
        print(f'{menu.menu} ({menu.price:,d}원)')
